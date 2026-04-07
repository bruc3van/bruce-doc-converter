import base64
import tempfile
import unittest
from datetime import date
from pathlib import Path
from unittest.mock import patch

import openpyxl
from docx import Document
from docx.enum.style import WD_STYLE_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx import Presentation
from pptx.util import Inches

from scripts.convert_document import (
    _extract_pdf_page_blocks,
    _postprocess_pdf_academic_sections,
    _render_docx_list_marker,
    batch_convert,
    convert_document,
)


class ConvertDocumentTests(unittest.TestCase):
    def test_convert_docx_skips_toc_paragraphs(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            docx_path = tmp_path / "toc.docx"
            output_dir = tmp_path / "out"

            document = Document()
            if "TOC 1" not in [style.name for style in document.styles]:
                document.styles.add_style("TOC 1", WD_STYLE_TYPE.PARAGRAPH)
            document.add_paragraph("目录", style="TOC Heading")
            document.add_paragraph("第一章\t1", style="TOC 1")
            document.add_paragraph("正文开始")
            document.save(docx_path)

            result = convert_document(str(docx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertNotIn("第一章\t1", result["markdown_content"])
            self.assertNotIn("目录", result["markdown_content"])
            self.assertIn("正文开始", result["markdown_content"])

    def test_render_docx_list_marker_preserves_multilevel_and_common_formats(self):
        numbering_state = {}
        levels = {
            0: {"start": 1, "num_fmt": "decimal", "lvl_text": "%1."},
            1: {"start": 1, "num_fmt": "decimal", "lvl_text": "%1.%2."},
            2: {"start": 1, "num_fmt": "decimal", "lvl_text": "%1.%2.%3."},
        }

        marker1 = _render_docx_list_marker({"ordered": True, "num_id": "n1", "level": 0, "levels": levels}, numbering_state)
        marker2 = _render_docx_list_marker({"ordered": True, "num_id": "n1", "level": 1, "levels": levels}, numbering_state)
        marker3 = _render_docx_list_marker({"ordered": True, "num_id": "n1", "level": 2, "levels": levels}, numbering_state)

        chinese_marker = _render_docx_list_marker(
            {"ordered": True, "num_id": "n2", "level": 0, "levels": {0: {"start": 1, "num_fmt": "chineseCounting", "lvl_text": "（%1）"}}},
            {},
        )
        circled_marker = _render_docx_list_marker(
            {"ordered": True, "num_id": "n3", "level": 0, "levels": {0: {"start": 1, "num_fmt": "decimalEnclosedCircle", "lvl_text": "%1"}}},
            {},
        )

        self.assertEqual("1.", marker1)
        self.assertEqual("1.1.", marker2)
        self.assertEqual("1.1.1.", marker3)
        self.assertEqual("（一）", chinese_marker)
        self.assertEqual("①", circled_marker)

    def test_convert_xlsx_keeps_merged_cells_as_single_value(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            xlsx_path = tmp_path / "merged.xlsx"
            output_dir = tmp_path / "out"

            workbook = openpyxl.Workbook()
            worksheet = workbook.active
            worksheet["A1"] = "Merged"
            worksheet.merge_cells("A1:C1")
            worksheet["A2"] = "v1"
            worksheet["B2"] = "v2"
            worksheet["C2"] = "v3"
            workbook.save(xlsx_path)
            workbook.close()

            result = convert_document(str(xlsx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertIn("| Merged |  |  |", result["markdown_content"])
            self.assertNotIn("| Merged | Merged | Merged |", result["markdown_content"])

    def test_convert_xlsx_supports_multi_headers_freeze_panes_and_multiple_tables(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            xlsx_path = tmp_path / "report.xlsx"
            output_dir = tmp_path / "out"

            workbook = openpyxl.Workbook()
            worksheet = workbook.active
            worksheet.freeze_panes = "A3"

            worksheet["A1"] = "Region"
            worksheet["B1"] = "Metrics"
            worksheet.merge_cells("B1:D1")
            worksheet["A2"] = "Name"
            worksheet["B2"] = "Date"
            worksheet["C2"] = "Rate"
            worksheet["D2"] = "Amount"
            worksheet["A3"] = "East"
            worksheet["B3"] = date(2024, 1, 2)
            worksheet["B3"].number_format = "yyyy-mm-dd"
            worksheet["C3"] = 0.125
            worksheet["C3"].number_format = "0.0%"
            worksheet["D3"] = 12345.6
            worksheet["D3"].number_format = "#,##0.00"

            worksheet["A6"] = 100
            worksheet["B6"] = 200
            worksheet["A7"] = 300
            worksheet["B7"] = 400

            workbook.save(xlsx_path)
            workbook.close()

            result = convert_document(str(xlsx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            markdown = result["markdown_content"]
            self.assertIn("### Table 1", markdown)
            self.assertIn("| Region / Name | Metrics / Date | Metrics / Rate | Metrics / Amount |", markdown)
            self.assertIn("| East | 2024-01-02 | 12.5% | 12,345.60 |", markdown)
            self.assertIn("### Table 2", markdown)
            self.assertIn("| Column 1 | Column 2 |", markdown)
            self.assertIn("| 100 | 200 |", markdown)
            self.assertIn("| 300 | 400 |", markdown)

    def test_convert_docx_vertical_merge_continuation_renders_blank_cell(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            docx_path = tmp_path / "vertical-merge.docx"
            output_dir = tmp_path / "out"

            document = Document()
            table = document.add_table(rows=3, cols=2)
            table.cell(0, 0).text = "Top"
            table.cell(1, 0).text = "Below"
            table.cell(0, 0).merge(table.cell(1, 0))
            table.cell(0, 1).text = "A"
            table.cell(1, 1).text = "B"
            table.cell(2, 0).text = "C"
            table.cell(2, 1).text = "D"
            document.save(docx_path)

            result = convert_document(str(docx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertIn("| Top Below | A |", result["markdown_content"])
            self.assertIn("|  | B |", result["markdown_content"])
            self.assertEqual(1, result["markdown_content"].count("| Top Below |"))

    def test_convert_pptx_allows_non_placeholder_textbox(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            pptx_path = tmp_path / "textbox.pptx"
            output_dir = tmp_path / "out"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            textbox.text_frame.text = "普通文本框"
            presentation.save(pptx_path)

            result = convert_document(str(pptx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertIn("普通文本框", result["markdown_content"])
            self.assertTrue(Path(result["output_path"]).exists())

    def test_convert_pptx_sorts_visual_order_and_splits_columns(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            pptx_path = tmp_path / "columns.pptx"
            output_dir = tmp_path / "out"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])

            right = slide.shapes.add_textbox(Inches(5.6), Inches(2.0), Inches(2.0), Inches(0.8))
            right.text_frame.text = "右侧内容"

            footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(2.0), Inches(0.4))
            footer.text_frame.text = "页脚"

            title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6.0), Inches(0.6))
            title.text_frame.text = "课程标题"

            left = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(2.0), Inches(0.8))
            left.text_frame.text = "左侧内容"

            presentation.save(pptx_path)

            result = convert_document(str(pptx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            markdown = result["markdown_content"]
            self.assertLess(markdown.index("### 课程标题"), markdown.index("#### Left Column"))
            self.assertIn("#### Left Column", markdown)
            self.assertIn("左侧内容", markdown)
            self.assertIn("#### Right Column", markdown)
            self.assertIn("右侧内容", markdown)
            self.assertIn("#### Footer", markdown)
            self.assertIn("页脚", markdown)

    def test_convert_pptx_extracts_chart_subtitle_and_notes(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            pptx_path = tmp_path / "chart-notes.pptx"
            output_dir = tmp_path / "out"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "季度回顾"
            slide.placeholders[1].text = "销售趋势"

            chart_data = CategoryChartData()
            chart_data.categories = ["Q1", "Q2"]
            chart_data.add_series("Sales", (12, 18))
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(1),
                Inches(2),
                Inches(5.5),
                Inches(3),
                chart_data,
            ).chart
            chart.has_title = True
            chart.chart_title.text_frame.text = "Revenue"
            slide.notes_slide.notes_text_frame.text = "讲解增长原因\n补充口径说明"
            presentation.save(pptx_path)

            result = convert_document(str(pptx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            markdown = result["markdown_content"]
            self.assertIn("### 季度回顾", markdown)
            self.assertIn("#### Subtitle", markdown)
            self.assertIn("销售趋势", markdown)
            self.assertIn("#### Visuals", markdown)
            self.assertIn("**Chart:** Revenue", markdown)
            self.assertIn("Series: Sales", markdown)
            self.assertIn("Categories: Q1, Q2", markdown)
            self.assertIn("### Notes", markdown)
            self.assertIn("讲解增长原因", markdown)

    def test_convert_pptx_groups_picture_caption(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            pptx_path = tmp_path / "picture-caption.pptx"
            output_dir = tmp_path / "out"
            image_path = tmp_path / "pixel.png"

            image_path.write_bytes(base64.b64decode(
                "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+jxM0AAAAASUVORK5CYII="
            ))

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_picture(str(image_path), Inches(1), Inches(1.5), Inches(2), Inches(2))
            caption = slide.shapes.add_textbox(Inches(1), Inches(3.65), Inches(2.8), Inches(0.5))
            caption.text_frame.text = "这是图片说明"
            presentation.save(pptx_path)

            result = convert_document(str(pptx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            markdown = result["markdown_content"]
            self.assertIn("#### Visuals", markdown)
            self.assertIn("**Image**", markdown)
            self.assertIn("Caption: 这是图片说明", markdown)

    def test_convert_docx_escapes_plain_markdown_syntax(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            docx_path = tmp_path / "plain.docx"
            output_dir = tmp_path / "out"

            document = Document()
            document.add_paragraph("1. 这是正文，不是列表")
            document.add_paragraph("# 这是正文，不是标题")
            document.save(docx_path)

            result = convert_document(str(docx_path), output_dir=str(output_dir))

            self.assertTrue(result["success"], result)
            self.assertIn("\\1. 这是正文，不是列表", result["markdown_content"])
            self.assertIn("\\# 这是正文，不是标题", result["markdown_content"])

    def test_convert_pdf_returns_clear_error_when_no_content_extracted(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            pdf_path = tmp_path / "empty.pdf"
            pdf_path.write_bytes(b"%PDF-1.4\n%%EOF\n")

            with patch("scripts.convert_document.check_dependencies", return_value=(True, None)):
                with patch("scripts.convert_document.convert_pdf", return_value=""):
                    result = convert_document(str(pdf_path))

            self.assertFalse(result["success"])
            self.assertIn("PDF 未提取到任何文本或表格", result["error"])

    def test_batch_convert_skips_generated_output_directories(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            root = Path(tmp_dir)
            (root / "source.docx").write_bytes(b"x")
            (root / "Markdown").mkdir()
            (root / "Markdown" / "generated.md").write_text("x", encoding="utf-8")
            (root / "Word").mkdir()
            (root / "Word" / "generated.docx").write_bytes(b"x")

            seen = []

            def _fake_convert(path, *_args, **_kwargs):
                seen.append(Path(path).relative_to(root).as_posix())
                return {"success": True}

            with patch("scripts.convert_document.convert_document", side_effect=_fake_convert):
                results = batch_convert(str(root), recursive=True)

            self.assertEqual(["source.docx"], seen)
            self.assertEqual(1, len(results))

    def test_extract_pdf_page_blocks_keeps_spanning_words_in_two_column_mode(self):
        words = [
            {"text": "FULLWIDTH", "x0": 10, "x1": 90, "top": 5, "bottom": 10, "upright": 1},
        ]
        for i in range(20):
            top = 20 + i * 5
            words.append({"text": f"L{i}", "x0": 5, "x1": 20, "top": top, "bottom": top + 4, "upright": 1})
            words.append({"text": f"R{i}", "x0": 80, "x1": 95, "top": top, "bottom": top + 4, "upright": 1})

        class FakePage:
            width = 100
            height = 200
            chars = []

            def __init__(self, page_words):
                self._words = page_words

            def filter(self, _predicate):
                return self

            def extract_words(self, **_kwargs):
                return list(self._words)

            def extract_text(self):
                return ""

        blocks = _extract_pdf_page_blocks(FakePage(words), tables=[])
        rendered = "".join(content for _, content in blocks)

        self.assertIn("FULLWIDTH", rendered)

    def test_postprocess_pdf_academic_sections_normalizes_common_blocks(self):
        content = """### Abstract

This is the abstract.

Keywords: alpha, beta；gamma

### References

[1] First ref

2. Second ref

### Appendix A

Proof details.
"""

        rendered = _postprocess_pdf_academic_sections(content)

        self.assertIn("## Abstract", rendered)
        self.assertIn("This is the abstract.", rendered)
        self.assertIn("## Keywords", rendered)
        self.assertIn("- alpha", rendered)
        self.assertIn("- beta", rendered)
        self.assertIn("- gamma", rendered)
        self.assertIn("## References", rendered)
        self.assertIn("1. First ref", rendered)
        self.assertIn("1. Second ref", rendered)
        self.assertIn("## Appendix A", rendered)


if __name__ == "__main__":
    unittest.main()
